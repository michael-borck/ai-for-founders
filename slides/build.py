#!/usr/bin/env python3
"""Build deck.pptx + deck.pdf from deck.html.

The HTML deck is the single source of truth. This renders each slide with
headless Chrome (via Playwright):

  deck.pptx — one full-bleed 16:9 image per slide, pixel-identical to the
              HTML, with invisible click-through rectangles over every link
              so QR/URL targets survive PowerPoint's own PDF export.
  deck.pdf  — Chrome print-to-PDF of the same file: vector text, selectable,
              links live, one slide per page (uses the deck's @page/@media
              print rules).
  deck-presenter.pptx — same slides plus speaker notes from
              facilitator/presenter-notes.md (## Slide N sections).
              Facilitator's copy — never send this one to the organiser.

Usage:  python3 slides/build.py
Needs:  playwright + Google Chrome, python-pptx.
"""

import re
import sys
import tempfile
from pathlib import Path

from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Emu, Inches

HERE = Path(__file__).resolve().parent
DECK = HERE / "deck.html"
PPTX_OUT = HERE / "deck.pptx"
PPTX_PRESENTER_OUT = HERE / "deck-presenter.pptx"
PDF_OUT = HERE / "deck.pdf"
NOTES_MD = HERE.parent / "facilitator" / "presenter-notes.md"

# 1280x720 CSS px matches Chrome's print layout for a 13.33in page, so the
# PPTX images fill the frame the same way the PDF pages do (at 1920 the
# 1000px content column + clamp()ed type float small in the middle).
VIEW_W, VIEW_H = 1280, 720           # CSS px used for layout + link rects
SLIDE_W_IN, SLIDE_H_IN = 40 / 3, 7.5  # 16:9 PowerPoint slide, inches


def render(tmpdir: Path):
    """Return (png_paths, link_rects_per_slide) and write PDF_OUT."""
    with sync_playwright() as p:
        browser = p.chromium.launch(channel="chrome")
        page = browser.new_page(
            viewport={"width": VIEW_W, "height": VIEW_H}, device_scale_factor=3
        )
        page.goto(DECK.as_uri())
        page.evaluate("document.fonts.ready")
        # handout images: no presenter chrome, no mid-fade captures
        page.add_style_tag(content=(
            ".nav,.counter,.brand,.progress,.hint{display:none!important}"
            "*{animation:none!important;transition:none!important}"
        ))
        n = page.evaluate("slides.length")

        pngs, links = [], []
        for i in range(n):
            page.evaluate(f"show({i})")
            png = tmpdir / f"slide-{i + 1:02d}.png"
            page.screenshot(path=str(png))
            pngs.append(png)
            links.append(page.evaluate(
                """[...document.querySelectorAll('.slide.active a')].map(a => {
                     const r = a.getBoundingClientRect();
                     return {href: a.href, x: r.x, y: r.y, w: r.width, h: r.height};
                   })"""
            ))

        page.evaluate("show(0)")
        page.pdf(path=str(PDF_OUT), prefer_css_page_size=True, print_background=True)
        browser.close()
    return pngs, links


def load_notes():
    """Map slide number -> notes text from presenter-notes.md, {} if absent."""
    if not NOTES_MD.exists():
        return {}
    parts = re.split(r"^## Slide (\d+)[^\n]*\n", NOTES_MD.read_text(), flags=re.M)
    return {int(parts[i]): parts[i + 1].strip() for i in range(1, len(parts), 2)}


def build_pptx(pngs, links, out, notes=None):
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)
    blank = prs.slide_layouts[6]

    def emu_x(px):
        return Emu(int(px / VIEW_W * prs.slide_width))

    def emu_y(px):
        return Emu(int(px / VIEW_H * prs.slide_height))

    for n, (png, slide_links) in enumerate(zip(pngs, links), start=1):
        slide = prs.slides.add_slide(blank)
        if notes and n in notes:
            slide.notes_slide.notes_text_frame.text = notes[n]
        slide.shapes.add_picture(
            str(png), 0, 0, width=prs.slide_width, height=prs.slide_height
        )
        for a in slide_links:
            shp = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                emu_x(a["x"]), emu_y(a["y"]), emu_x(a["w"]), emu_y(a["h"]),
            )
            shp.fill.background()
            shp.line.fill.background()
            shp.shadow.inherit = False
            shp.click_action.hyperlink.address = a["href"]

    prs.save(out)


def main():
    with tempfile.TemporaryDirectory() as td:
        pngs, links = render(Path(td))
        build_pptx(pngs, links, PPTX_OUT)
        notes = load_notes()
        if notes:
            build_pptx(pngs, links, PPTX_PRESENTER_OUT, notes)
        made = [PPTX_OUT, PDF_OUT] + ([PPTX_PRESENTER_OUT] if notes else [])
        print(f"{len(pngs)} slides -> " + ", ".join(
            f"{p.name} ({p.stat().st_size // 1024} KB)" for p in made))
    return 0


if __name__ == "__main__":
    sys.exit(main())
